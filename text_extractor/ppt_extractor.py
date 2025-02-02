from pptx import Presentation


def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    full_text = ''

    for i, slide in enumerate(prs.slides):
        # Extract slide header (title)
        header = ''
        if slide.shapes.title:
            header = slide.shapes.title.text.strip()

        # Extract slide content (other text)
        content = ''
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape != slide.shapes.title:
                content += shape.text.strip() + '\n'

        # Combine header and content for the current slide
        slide_text = f"Slide {i + 1} Header: {header}\nSlide {i + 1} Content:\n{content.strip()}"
        full_text += slide_text + '\n---\n'  # Add delimiter between slides

    return full_text.strip()  # Remove the last delimiter
